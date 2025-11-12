from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Aurkezpen berri bat sortu
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def gehitu_titulua(slide, testua, subtitle=None):
    """Titulu slide bat gehitu"""
    title = slide.shapes.title
    title.text = testua
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def gehitu_edukia(slide, titulua, edukia):
    """Eduki slide bat gehitu"""
    title = slide.shapes.title
    title.text = titulua
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = edukia
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.level = 0

# SLIDE 1: Titulua
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
gehitu_titulua(slide, "🎯 IKUSMENEKO SISTEMAK", 
              "Lehen Pausuak Computer Vision-en\n\nAdimen Artifiziala - Oinarrizko Maila")

# SLIDE 2: Zer da Computer Vision?
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "🤖 ZER DA COMPUTER VISION?", 
"""Computer Vision edo Ikusmen Artifiziala ordenagailuei irudiak eta bideoak "ikusteko" eta ulertzeko gaitasuna ematen dien teknologia da.

Giza ikusmena imitatzen saiatzen da!

• Irudiak analizatu
• Objektuak detektatu
• Patroi eta ereduak ezagutu
• Erabakiak hartu informazio bisualean oinarrituta""")

# SLIDE 3: Irudi Digitala - Oinarriak
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "📸 IRUDI DIGITALA", 
"""Irudi bat PIXELAK-ek osatzen dute.

Pixel = Picture Element
(Irudiaren elementu txikiena)

Pixel bakoitzak KOLORE bat du!

• Pixelak: Irudiaren oinarrizko unitateak
• Bereizmena: Pixel kopurua (adib: 1920x1080)
• Kolore sakonera: Informazio kantitatea pixel bakoitzeko""")

# SLIDE 4: RGB Kolore Sistema
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "🎨 RGB: Kolore Guztiak 3 Koloreetatik!", 
"""R = Red (Gorria)
G = Green (Berdea)
B = Blue (Urdina)

Kolore bakoitzak 0-255 bitarteko balioa du

ADIBIDEAK:
• Zuria = R:255, G:255, B:255
• Beltza = R:0, G:0, B:0
• Horia = R:255, G:255, B:0
• Gorria = R:255, G:0, B:0""")

# SLIDE 5: Irudi Motak
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "🖼️ IRUDI MOTA DESBERDINAK", 
"""1. Kolore Irudiak (RGB)
   • 3 kanal
   • Kolore osoa

2. Eskala Griseko Irudiak
   • 1 kanal
   • Zuri-beltzetik

3. Irudi Binariak
   • 0 edo 1 soilik
   • Beltz edo zuria""")

# SLIDE 6: Aplikazio Errealak
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "🌍 APLIKAZIO ERREALAK", 
"""Ikusmeneko sistemak EDONON daude!

🔒 Segurtasuna: Aurpegi-detekzioa, matrikula-irakurketa
🏥 Osasuna: RX irudien analisia, diagnostiko laguntza
🏭 Industria: Kalitateko kontrola, robotika
📱 Teknologia: Face ID, filtroak, QR kodeak
🚗 Garraioa: Ibilgailu autonomoak
🌾 Nekazaritza: Uzta monitorizazioa, droneak""")

# SLIDE 7: Python + OpenCV
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "💻 LEHEN PAUSUAK PROGRAMAZIOAN", 
"""OpenCV = Open Source Computer Vision Library
Tresnarik ezagunena Computer Vision-erako!

Python lengoaia erabiliko dugu (erraza!)

KODE ADIBIDEA:
import cv2

# Irudia kargatu
irudia = cv2.imread('nire_irudia.jpg')

# Irudia bistaratu
cv2.imshow('Nire Irudia', irudia)
cv2.waitKey(0)""")

# SLIDE 8: Praktika Sinplea
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "🎯 PRAKTIKA 1: Zure lehen programa!", 
"""PAUSUAK:

1. Python instalatu (3.8 edo berriagoa)

2. OpenCV instalatu:
   pip install opencv-python

3. Deskargatu test irudia

4. Kopiatu kodea eta exekutatu

5. Emaitza ikusi!""")

# SLIDE 9: Laburpena
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "📝 LABURPENA", 
"""IKUSMENEKO SISTEMAK

OINARRIAK:
• Pixelak - Irudien oinarrizko unitateak
• RGB - Kolore sistema (3 kanal)
• Formatuak - RGB, Grisa, Binarioa

APLIKAZIOAK:
• Segurtasuna, Osasuna, Industria
• Teknologia, Garraioa, Nekazaritza

TRESNAK:
• Python programazio lengoaia
• OpenCV liburutegia""")

# SLIDE 10: Hurrengo Pausuak
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
gehitu_edukia(slide, "🚀 BIKAIN! Oinarriak ikasi dituzu!", 
"""Orain prest zaude hurrengo atalera pasatzeko:

➡️ 2.B ATAZA: Ezagutzan Sakontzea

• Deep Learning
• CNN (Convolutional Neural Networks)
• Transfer Learning
• Proiektu aurreratuak

Jarraitu ikasten! 💪""")

# Gorde aurkezpena
fitxategia = "Ikusmeneko_Sistemak_Oinarrizko_Maila.pptx"
prs.save(fitxategia)
print(f"✅ Aurkezpena sortuta: {fitxategia}")
print(f"📊 Diapositibak: {len(prs.slides)}")
