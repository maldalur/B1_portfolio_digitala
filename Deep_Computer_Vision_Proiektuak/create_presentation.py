"""
Deep Computer Vision Proiektuaren Aurkezpena
PowerPoint sortzeko script-a
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def sortu_aurkezpena():
    """PowerPoint aurkezpena sortu"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # --- SLIDE 1: Azala ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Fondo kolorea
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 35, 126)  # Urdin iluna
    
    # Titulua
    txbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txbox.text_frame
    tf.text = "🚀 Deep Computer Vision"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Azpititulua
    txbox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    tf2 = txbox2.text_frame
    tf2.text = "Ikastaro Aurreratua - Portfolio Digitala"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(255, 193, 7)  # Horia
    p2.alignment = PP_ALIGN.CENTER
    
    # Egilea
    txbox3 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf3 = txbox3.text_frame
    tf3.text = "Mikel Aldalur Corta\nInstituto BIRT - 2025"
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(255, 255, 255)
    p3.alignment = PP_ALIGN.CENTER
    
    # --- SLIDE 2: Helburua ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎯 Ikastaroaren Helburua"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Maila aurreratuko ikasleei Deep Learning eta Computer Vision teknikak irakastea"
    p.font.size = Pt(20)
    p.space_after = Pt(20)
    
    bullets = [
        "CNN ereduak zerotik eraikitzen ikasi",
        "Transfer Learning teknikak aplikatu",
        "Objektu detekzioa denbora errealean",
        "Aurpegi ezagutza sistemak garatu",
        "Proiektu errealetan aplikatzeko gaitasuna lortu"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    # --- SLIDE 3: Ikastaroaren Egitura ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📚 Ikastaroaren Egitura"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    atalak = [
        ("01 - CNN Oinarriak", "2 notebooks, 2-3h"),
        ("02 - Transfer Learning", "2 notebooks, 2-3h"),
        ("03 - Objektu Detekzioa", "2 notebooks, 2-3h"),
        ("04 - Aurpegi Ezagutza", "2 notebooks, 2-3h"),
        ("05 - Proiektu Finala", "1 notebook, 2-3h")
    ]
    
    for atala, info in atalak:
        p = tf.add_paragraph()
        p.text = f"{atala}\n     {info}"
        p.font.size = Pt(18)
        p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "\nGUZTIRA: 9 notebooks, 10-15 ordu"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(76, 175, 80)
    
    # --- SLIDE 4: Atala 1 - CNN Oinarriak ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📘 Atala 1: CNN Oinarriak"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "01_CNN_Sarrera.ipynb"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(10)
    
    bullets1 = [
        "CNN kontzeptuaren sarrera",
        "MLP vs CNN konparaketa",
        "MNIST dataset-arekin praktika",
        "Feature maps bistaratzea"
    ]
    
    for bullet in bullets1:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\n02_Lehen_CNN_Eredua.ipynb"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(10)
    
    bullets2 = [
        "CNN eredua zerotik eraikitzen",
        "Callbacks eta optimizazioa",
        "Confusion matrix eta metrikak"
    ]
    
    for bullet in bullets2:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    # --- SLIDE 5: Atala 2 - Transfer Learning ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📗 Atala 2: Transfer Learning"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "01_Transfer_Learning_Sarrera.ipynb"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(10)
    
    bullets1 = [
        "Transfer Learning kontzeptua",
        "VGG16, ResNet50, MobileNetV2",
        "Feature Extraction vs Fine-Tuning"
    ]
    
    for bullet in bullets1:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\n02_Cats_vs_Dogs_Klasifikazioa.ipynb"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(10)
    
    bullets2 = [
        "Katu vs Txakur klasifikazioa",
        "Data Augmentation aplikatu",
        "Custom dataset-arekin lan egin"
    ]
    
    for bullet in bullets2:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    # --- SLIDE 6: Atala 3 - Objektu Detekzioa ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📙 Atala 3: Objektu Detekzioa"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "01_Objektu_Detekzioa_Sarrera.ipynb"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(10)
    
    bullets1 = [
        "Objektu detekzioaren oinarriak",
        "Bounding boxes, IoU, NMS",
        "YOLO, R-CNN, SSD algoritmoak"
    ]
    
    for bullet in bullets1:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\n02_YOLO_Praktika.ipynb"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(10)
    
    bullets2 = [
        "YOLOv8 erabilera praktikoa",
        "Irudietan eta bideoetan detektatu",
        "Real-time detekzioa webcam-ekin"
    ]
    
    for bullet in bullets2:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    # --- SLIDE 7: Atala 4 - Aurpegi Ezagutza ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📕 Atala 4: Aurpegi Ezagutza"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "01_Aurpegi_Ezagutza_Sarrera.ipynb"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(10)
    
    bullets1 = [
        "Face Detection vs Recognition",
        "Face Embeddings (128D)",
        "Verification vs Identification"
    ]
    
    for bullet in bullets1:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\n02_Face_Recognition_Praktika.ipynb"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(10)
    
    bullets2 = [
        "face_recognition liburutegia",
        "Face landmarks (68 puntu)",
        "Real-time recognition webcam-ekin"
    ]
    
    for bullet in bullets2:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    # --- SLIDE 8: Atala 5 - Proiektu Finala ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📓 Atala 5: Proiektu Finala"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "01_Proiektu_Finala_CIFAR10.ipynb"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(15)
    
    bullets = [
        "CIFAR-10 dataset osoa (60,000 irudi)",
        "Custom CNN eredua diseinatu",
        "Data Augmentation estrategia",
        "Ebaluazio integrala",
        "Eredua gorde eta exportatu",
        "Ondorioak eta hurrengo pausoak"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # --- SLIDE 9: Teknologiak ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🛠️ Teknologiak eta Tresnak"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Software"
    p.font.size = Pt(22)
    p.font.bold = True
    
    bullets1 = ["Python 3.8+", "Jupyter Notebook", "Google Colab (GPU)"]
    for bullet in bullets1:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\nLiburutegiak"
    p.font.size = Pt(22)
    p.font.bold = True
    
    bullets2 = [
        "TensorFlow / Keras",
        "OpenCV",
        "NumPy, Matplotlib, Seaborn",
        "scikit-learn",
        "Ultralytics (YOLO)",
        "face_recognition (dlib)"
    ]
    
    for bullet in bullets2:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(16)
    
    # --- SLIDE 10: Lorpen Esperatuak ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🏆 Zer Lortu Dezakezu"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    lorpenak = [
        "CNN ereduak zerotik eraikitzen",
        "Transfer Learning aplikatzen",
        "Objektuak denbora errealean detektatzen",
        "Aurpegi-ezagutza sistemak garatzen",
        "Computer Vision proiektu konplexuak sortzen",
        "Ereduak produkziora eramaten",
        "Portfolio digital profesionala izaten"
    ]
    
    for lorpen in lorpenak:
        p = tf.add_paragraph()
        p.text = f"✅ {lorpen}"
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # --- SLIDE 11: Hurrengo Pausoak ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🚀 Hurrengo Pausoak"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    pausoak = [
        ("1. Dataset Propioa", "Zure datu propioekin entrenatu"),
        ("2. API Sortu", "Flask/FastAPI REST API"),
        ("3. Optimizatu", "TensorFlow Lite, ONNX"),
        ("4. Deploy", "AWS, Google Cloud, Azure"),
        ("5. Mobil App", "TensorFlow Lite mugikorrean"),
        ("6. Portfolio", "GitHub-en argitaratu proiektuak")
    ]
    
    for pauso, deskripzioa in pausoak:
        p = tf.add_paragraph()
        p.text = f"{pauso}\n     {deskripzioa}"
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # --- SLIDE 12: Esker Oneko Orria ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo kolorea
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(76, 175, 80)  # Berde
    
    # Eskerrik asko
    txbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txbox.text_frame
    tf.text = "Eskerrik asko! 🎉\n\nGalderak?"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Kontaktua
    txbox2 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf2 = txbox2.text_frame
    tf2.text = "Mikel Aldalur Corta\nmikel.aldalur@birt.eus\nGitHub: github.com/maldalur"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    # Gorde
    output_file = "Deep_Computer_Vision_Aurkezpena.pptx"
    prs.save(output_file)
    print(f"✅ Aurkezpena sortuta: {output_file}")
    return output_file

if __name__ == "__main__":
    try:
        sortu_aurkezpena()
    except Exception as e:
        print(f"❌ Errorea: {e}")
        print("\n📦 python-pptx instalatu behar duzu:")
        print("   pip install python-pptx")
